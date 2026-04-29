# rina_media.py
import base64
import os
from PIL import Image
from io import BytesIO

# 导入新增的文档解析库
import fitz  # PyMuPDF 处理 PDF
import docx  # 处理 Word
import pptx  # 处理 PPT
import pandas as pd  # 处理 Excel/CSV


def extract_text_from_pdf(bytes_data):
    """从 PDF 提取文本"""
    text = ""
    # fitz 支持直接从内存字节流读取
    with fitz.open(stream=bytes_data, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text


def extract_text_from_docx(bytes_data):
    """从 Word 提取文本"""
    doc = docx.Document(BytesIO(bytes_data))
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_from_pptx(bytes_data):
    """从 PPT 提取文本"""
    prs = pptx.Presentation(BytesIO(bytes_data))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def extract_text_from_excel(bytes_data, filename):
    """从表格提取数据并转换为 Markdown 格式"""
    if filename.endswith('.csv'):
        df = pd.read_csv(BytesIO(bytes_data))
    else:
        df = pd.read_excel(BytesIO(bytes_data))
    # 转换为 Markdown 表格格式，大模型阅读体验最好
    return df.to_markdown(index=False)


def process_uploaded_file(uploaded_file):
    """处理上传文件载荷，转换为模型可读格式"""
    if uploaded_file is None:
        return None

    file_type = uploaded_file.type
    filename = uploaded_file.name.lower()
    bytes_data = uploaded_file.getvalue()

    # 1. 图像处理保持不变
    if file_type.startswith("image/"):
        base64_image = base64.b64encode(bytes_data).decode('utf-8')
        return {
            "type": "image_url",
            "image_url": {"url": f"data:{file_type};base64,{base64_image}"}
        }

    # 2. 文本与富文本文档处理
    else:
        try:
            parsed_text = ""

            # 根据后缀名进行分发解析
            if filename.endswith(".pdf"):
                parsed_text = extract_text_from_pdf(bytes_data)
            elif filename.endswith((".docx", ".doc")):
                parsed_text = extract_text_from_docx(bytes_data)
            elif filename.endswith((".pptx", ".ppt")):
                parsed_text = extract_text_from_pptx(bytes_data)
            elif filename.endswith((".xlsx", ".xls", ".csv")):
                parsed_text = extract_text_from_excel(bytes_data, filename)
            else:
                # 兜底方案：如果是纯文本代码文件 (.py, .c, .txt, .md 等)
                parsed_text = bytes_data.decode("utf-8")

            # 安全校验：防止扫描版图片 PDF 提取出空字符串
            if not parsed_text.strip():
                parsed_text = "[文档解析警告：未提取到有效文本，这可能是一个纯图片扫描版文档或空文档。]"

            return {
                "type": "text",
                "text": f"\n--- [载入文件: {uploaded_file.name}] ---\n{parsed_text}\n--- [文件结束] ---\n"
            }

        except Exception as e:
            return {"type": "text", "text": f"\n[文件读取物理故障: {str(e)}]\n"}


def save_input_image(uploaded_file):
    """将输入流中的图像硬写入本地 SSD"""
    if uploaded_file is None:
        return None

    if not os.path.exists("img"):
        os.makedirs("img")

    img_path = f"img/input_{os.urandom(4).hex()}.png"
    bytes_data = uploaded_file.getvalue()
    image = Image.open(BytesIO(bytes_data))
    image.save(img_path)
    return img_path